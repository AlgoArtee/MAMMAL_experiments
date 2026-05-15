import sys
from pathlib import Path

sys.path.insert(0, str(Path(".deps").resolve()))

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path("MAMMAL_Kernaussagen_Praesentation.pptx")

WIDE_W = 13.333
WIDE_H = 7.5

NAVY = RGBColor(18, 35, 52)
TEAL = RGBColor(0, 122, 130)
GREEN = RGBColor(52, 151, 92)
AMBER = RGBColor(221, 151, 50)
RED = RGBColor(184, 72, 67)
GRAY = RGBColor(91, 105, 120)
LIGHT = RGBColor(244, 247, 250)
MID = RGBColor(218, 226, 235)
WHITE = RGBColor(255, 255, 255)


def add_bg(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def text_box(slide, x, y, w, h, text, size=20, bold=False, color=NAVY, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    run = p.runs[0]
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def title(slide, text, subtitle=None):
    text_box(slide, 0.55, 0.28, 10.9, 0.5, text, 24, True, NAVY)
    if subtitle:
        text_box(slide, 0.58, 0.83, 10.9, 0.35, subtitle, 10.5, False, GRAY)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(1.18), Inches(1.15), Inches(0.045))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()


def footer(slide, idx):
    text_box(
        slide,
        0.58,
        7.13,
        9.3,
        0.22,
        "Quelle: Shoshan et al., npj Drug Discovery 2026, DOI 10.1038/s44386-026-00047-4",
        7.5,
        False,
        GRAY,
    )
    text_box(slide, 12.3, 7.13, 0.45, 0.22, str(idx), 8, False, GRAY, PP_ALIGN.RIGHT)


def pill(slide, x, y, w, h, label, color=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = label
    r = p.runs[0]
    r.font.name = "Aptos"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = WHITE
    return shape


def bullet_list(slide, x, y, w, h, items, size=15, color=NAVY):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(5)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p._p.get_or_add_pPr().set("marL", "285750")
        p._p.get_or_add_pPr().set("indent", "-171450")
    return box


def card(slide, x, y, w, h, heading, body, accent=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = MID
    shape.line.width = Pt(0.8)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    text_box(slide, x + 0.22, y + 0.18, w - 0.35, 0.28, heading, 13, True, NAVY)
    text_box(slide, x + 0.22, y + 0.58, w - 0.35, h - 0.78, body, 11, False, GRAY)


def add_table(slide, x, y, w, h, rows, col_widths=None, font_size=9.5):
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    if col_widths:
        for i, width in enumerate(col_widths):
            table.columns[i].width = Inches(width)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r == 0 else (LIGHT if r % 2 else WHITE)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(font_size if r else font_size - 0.5)
                p.font.bold = r == 0
                p.font.color.rgb = WHITE if r == 0 else NAVY
    return table


def metric_card(slide, x, y, value, label, color=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.1), Inches(1.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    text_box(slide, x + 0.12, y + 0.15, 1.85, 0.38, value, 21, True, WHITE, PP_ALIGN.CENTER)
    text_box(slide, x + 0.14, y + 0.61, 1.82, 0.28, label, 8.5, False, WHITE, PP_ALIGN.CENTER)


prs = Presentation()
prs.slide_width = Inches(WIDE_W)
prs.slide_height = Inches(WIDE_H)
blank = prs.slide_layouts[6]


# 1
slide = prs.slides.add_slide(blank)
add_bg(slide, NAVY)
text_box(slide, 0.65, 0.55, 11.7, 0.7, "MAMMAL", 36, True, WHITE)
text_box(
    slide,
    0.68,
    1.34,
    10.6,
    0.7,
    "Molecular Aligned Multi-Modal Architecture and Language for biomedical discovery",
    20,
    False,
    RGBColor(205, 230, 235),
)
text_box(
    slide,
    0.7,
    2.35,
    8.7,
    1.1,
    "Kernaussage: Ein multimodales Foundation Model kann Protein-, Antikoerper-, Molekuel- und Genexpressionsdaten in einer gemeinsamen Sequenzsprache nutzen und erreicht in 9 von 11 Drug-Discovery-Benchmarks neue SOTA-Werte.",
    24,
    True,
    WHITE,
)
metric_card(slide, 0.7, 4.45, "2 Mrd.", "Pretraining-Samples", TEAL)
metric_card(slide, 3.05, 4.45, "11", "Benchmarks", GREEN)
metric_card(slide, 5.4, 4.45, "9/11", "SOTA-Ergebnisse", AMBER)
metric_card(slide, 7.75, 4.45, "5/7", "besser als AF3", RED)
text_box(slide, 0.7, 6.85, 8.6, 0.22, "Erstellt aus s44386-026-00047-4.pdf; PDF unveraendert.", 8, False, RGBColor(196, 208, 218))


# 2
slide = prs.slides.add_slide(blank)
add_bg(slide)
title(slide, "Warum das relevant ist", "Drug discovery braucht Modelle, die Modalitaeten und Aufgaben verbinden.")
footer(slide, 2)
card(slide, 0.7, 1.65, 3.75, 1.45, "Ausgangsproblem", "Klassische Pipelines sind teuer, langsam und leiden unter hoher klinischer Ausfallrate; viele Datenarten bleiben getrennt modelliert.", RED)
card(slide, 4.8, 1.65, 3.75, 1.45, "Modell-Luecke", "Molekuele, Proteine, Antikoerper und Genexpression werden oft mit spezialisierten Architekturen statt gemeinsamem Kontext verarbeitet.", AMBER)
card(slide, 8.9, 1.65, 3.75, 1.45, "Zielbild", "Ein Modell soll Vorhersage, Regression und Generierung ueber mehrere Stufen der Wirkstoffentwicklung hinweg unterstuetzen.", TEAL)
bullet_list(
    slide,
    1.05,
    3.72,
    11.2,
    1.75,
    [
        "Die Arbeit positioniert MAMMAL als integratives Modell fuer Zielidentifikation, Wirkstoffwirkung, Sicherheit, Bindung und Antikoerperdesign.",
        "Die Evaluation ist bewusst benchmark-orientiert: vorhandene Splits, etablierte Metriken und Vergleich gegen spezialisierte SOTA-Modelle.",
        "Die praktische Relevanz wird zusaetzlich durch eine experimentell validierte Wirkstoff-Rangfolge und AF3-Vergleiche gestuetzt.",
    ],
    15,
)


# 3
slide = prs.slides.add_slide(blank)
add_bg(slide)
title(slide, "Loesung: Eine gemeinsame Sequenzsprache", "MAMMAL uebersetzt heterogene biologische Entitaeten in strukturierte Prompts.")
footer(slide, 3)
for x, label, color in [(0.65, "Small Molecules\nSMILES", TEAL), (3.25, "Proteine\nAminosaeuren", GREEN), (5.85, "Antikoerper\nHeavy/Light Chains", AMBER), (8.45, "Genexpression\nsortierte Gene", RED)]:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.65), Inches(2.15), Inches(0.95))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    text_box(slide, x + 0.1, 1.82, 1.95, 0.45, label, 12, True, WHITE, PP_ALIGN.CENTER)
pill(slide, 4.6, 3.05, 3.25, 0.55, "Modularer Tokenizer + strukturierte Prompts", NAVY)
pill(slide, 4.83, 4.05, 2.8, 0.55, "Transformer mit Encoder-only und Encoder-Decoder", TEAL)
bullet_list(
    slide,
    0.9,
    5.05,
    11.8,
    1.2,
    [
        "Zahlen werden nicht als Textziffern zerlegt, sondern ueber Projektionen direkt in den Embedding-Raum integriert.",
        "Dasselbe Framework kann Klassifikation, Regression und Generierung bedienen, statt fuer jede Aufgabe ein eigenes Modell zu bauen.",
    ],
    14,
)


# 4
slide = prs.slides.add_slide(blank)
add_bg(slide)
title(slide, "Beleg 1: Breite Benchmark-Abdeckung", "Leistung ueber elf Aufgaben entlang der Wirkstoffpipeline.")
footer(slide, 4)
rows = [
    ["Benchmark", "Domain", "Metrik", "SOTA", "MAMMAL", "Imp."],
    ["Cell type", "GE", "F1", "0.710", "0.763", "+7.5%"],
    ["BBBP", "SM", "AUROC", "0.937", "0.957", "+2.2%"],
    ["ClinTox", "SM", "AUROC", "0.948", "0.986", "+4.0%"],
    ["Cancer Drug 1", "GE+SM", "Pearson", "0.887", "0.917", "+3.4%"],
    ["Cancer Drug 2", "GE+SM", "Pearson", "0.900", "0.931", "+3.4%"],
    ["Ab Infilling", "Protein", "CDRH3-AAR", "0.375", "0.446", "+19.0%"],
    ["PPI ddG", "Protein", "Pearson", "0.663", "0.852", "+28.5%"],
    ["DTI", "Prot.+SM", "NRMSE", "0.942", "0.906", "+3.8%"],
]
add_table(slide, 0.62, 1.45, 7.65, 4.95, rows, [1.65, 0.75, 1.0, 0.75, 0.9, 0.7], 8.8)
chart_data = CategoryChartData()
chart_data.categories = ["Cell", "ClinTox", "Ab Inf.", "PPI", "DTI"]
chart_data.add_series("Relative Verbesserung", [7.5, 4.0, 19.0, 28.5, 3.8])
chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(8.65), Inches(1.7), Inches(3.9), Inches(3.6), chart_data).chart
chart.has_legend = False
chart.value_axis.has_major_gridlines = True
chart.value_axis.tick_labels.font.size = Pt(8)
chart.category_axis.tick_labels.font.size = Pt(8)
chart.plots[0].series[0].format.fill.solid()
chart.plots[0].series[0].format.fill.fore_color.rgb = TEAL
text_box(slide, 8.7, 5.58, 3.65, 0.62, "Wichtig: Zwei Aufgaben sind nur kompetitiv bzw. knapp verbessert; der SOTA-Claim gilt innerhalb der gewaehlten oeffentlichen Benchmarks.", 10.5, False, GRAY)


# 5
slide = prs.slides.add_slide(blank)
add_bg(slide)
title(slide, "Beleg 2: Vergleich gegen AlphaFold 3", "MAMMAL wird als Sequenzmodell gegen AF3-Konfidenzscores fuer Binder/Nicht-Binder getestet.")
footer(slide, 5)
rows = [
    ["Target", "n", "MAMMAL", "AF3", "Delta"],
    ["HER2 ECD", "60", "0.93", "0.45", "+0.42"],
    ["Albumin", "28", "0.91", "0.59", "+0.32"],
    ["CD206", "37", "1.00", "0.59", "+0.41"],
    ["EGFR", "28", "0.94", "0.49", "+0.45"],
    ["TBG", "32", "0.63", "1.00", "-0.37"],
    ["TNF alpha", "34", "0.86", "0.87", "-0.01"],
    ["VWF", "34", "0.83", "0.32", "+0.51"],
]
add_table(slide, 0.72, 1.45, 5.35, 4.25, rows, [1.3, 0.5, 0.85, 0.65, 0.7], 9.5)
chart_data = CategoryChartData()
chart_data.categories = ["HER2", "Alb.", "CD206", "EGFR", "TBG", "TNF", "VWF"]
chart_data.add_series("MAMMAL", [0.93, 0.91, 1.00, 0.94, 0.63, 0.86, 0.83])
chart_data.add_series("AF3", [0.45, 0.59, 0.59, 0.49, 1.00, 0.87, 0.32])
chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(6.45), Inches(1.48), Inches(6.05), Inches(4.25), chart_data).chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.value_axis.maximum_scale = 1.0
chart.value_axis.tick_labels.font.size = Pt(8)
chart.category_axis.tick_labels.font.size = Pt(8)
chart.plots[0].series[0].format.fill.solid()
chart.plots[0].series[0].format.fill.fore_color.rgb = TEAL
chart.plots[0].series[1].format.fill.solid()
chart.plots[0].series[1].format.fill.fore_color.rgb = AMBER
bullet_list(slide, 0.8, 6.05, 11.7, 0.75, ["Interpretation: AF3 liefert 3D-Hypothesen; MAMMAL liefert diskriminative Bindungswahrscheinlichkeiten. Die Studie vergleicht daher Entscheidungsnutzen, nicht identische Modellziele."], 12.5)


# 6
slide = prs.slides.add_slide(blank)
add_bg(slide)
title(slide, "Beleg 3: Wirkstoffantwort mit Nasslabor-Anker", "Vier nicht im Training enthaltene Wirkstoffe wurden in Zelllinien experimentell geprueft.")
footer(slide, 6)
for i, (drug, note, color) in enumerate([
    ("Carfilzomib", "staerkste Potenz", GREEN),
    ("Nintedanib", "Rang 2", TEAL),
    ("Infigratinib", "Rang 3", AMBER),
    ("Vemurafenib", "geringste Potenz", RED),
]):
    x = 0.95 + i * 3.0
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(1.8), Inches(1.8), Inches(1.05))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    text_box(slide, x - 0.2, 3.05, 2.25, 0.28, drug, 12.5, True, NAVY, PP_ALIGN.CENTER)
    text_box(slide, x - 0.2, 3.38, 2.25, 0.25, note, 9.5, False, GRAY, PP_ALIGN.CENTER)
bullet_list(
    slide,
    1.0,
    4.45,
    11.2,
    1.25,
    [
        "MAMMAL reproduzierte die experimentelle Potenz-Rangfolge fuer die getesteten Zelllinien.",
        "Auf 805 GDSC-Zelllinien blieb die relative Ordnung in ca. 90-95% der Faelle erhalten.",
        "Carfilzomib ist bisher fuer haematologische Malignitaeten zugelassen; das Signal in soliden Tumorzelllinien begruendet Follow-up, nicht unmittelbare klinische Schlussfolgerung.",
    ],
    14,
)


# 7
slide = prs.slides.add_slide(blank)
add_bg(slide)
title(slide, "Entscheidungen aus der Arbeit", "Welche Architektur- und Forschungsentscheidungen tragen die Ergebnisse?")
footer(slide, 7)
card(slide, 0.72, 1.55, 3.8, 1.55, "1. Multi-Domain statt Einzeldomain", "Gemeinsame Syntax fuer SMILES, Proteinsequenzen, Antikoerperketten und Genexpression.", TEAL)
card(slide, 4.82, 1.55, 3.8, 1.55, "2. Zahlen nativ modellieren", "Kontinuierliche Werte werden projiziert; das ist zentral fuer IC50, pKD, ddG und andere Skalen.", GREEN)
card(slide, 8.92, 1.55, 3.8, 1.55, "3. Zwei Betriebsarten", "Encoder-only fuer Klassifikation/Regression, Encoder-Decoder fuer Generierung und Infilling.", AMBER)
card(slide, 0.72, 3.65, 3.8, 1.55, "4. SOTA nur mit fairen Splits", "Benchmark-Auswahl priorisiert oeffentliche Daten, definierte Splits und etablierte Metriken.", NAVY)
card(slide, 4.82, 3.65, 3.8, 1.55, "5. AF3 gezielt einsetzen", "AF3 bleibt fuer Strukturhypothesen wertvoll; MAMMAL ist in Binder-Klassifikation oft entscheidungsstaerker.", RED)
card(slide, 8.92, 3.65, 3.8, 1.55, "6. Offen bereitstellen", "Code, Modellgewichte und Tokenizer sind als Grundlage fuer Reproduktion und weitere Studien verfuegbar.", TEAL)


# 8
slide = prs.slides.add_slide(blank)
add_bg(slide)
title(slide, "Was man nicht ueberinterpretieren sollte", "Die Evidenz ist stark, aber nicht frei von Grenzen.")
footer(slide, 8)
bullet_list(
    slide,
    0.9,
    1.55,
    5.75,
    4.5,
    [
        "Die SOTA-Aussage haengt an den berichteten oeffentlichen Benchmarks und deren Splits.",
        "AF3-Vergleiche nutzen Konfidenzscores als Proxy; AF3 ist kein explizit trainierter Binder-Klassifikator.",
        "Einige AF3-Tests sind klein, z. B. HER2 mit 60 Beispielen und mehrere Targets unter 40 Beispielen.",
        "Wirkstoffantworten aus Zelllinien und CellTiter-Glo sind fruehe Evidenz; klinische Uebertragbarkeit bleibt offen.",
    ],
    15,
)
card(slide, 7.2, 1.65, 4.75, 1.45, "Implikation", "MAMMAL ist eine robuste Forschungsplattform fuer Priorisierung und Hypothesengenerierung, aber kein Ersatz fuer prospektive biologische Validierung.", RED)
card(slide, 7.2, 3.55, 4.75, 1.45, "Governance", "Fuer produktiven Einsatz braucht es dataset-spezifische Validierung, Unsicherheitskommunikation und klare Grenzen je Task.", AMBER)


# 9
slide = prs.slides.add_slide(blank)
add_bg(slide)
title(slide, "Naechste Schritte", "Pragmatischer Weg von Paper zu nutzbarer Entscheidungsunterstuetzung.")
footer(slide, 9)
steps = [
    ("1", "Reproduzieren", "Code/Weights laden, Zielbenchmarks lokal nachbauen, Metriken gegen Paper-Werte pruefen."),
    ("2", "Use Case waehlen", "Ein klarer Arbeitsfluss: z. B. DTI-Priorisierung, Antikoerper-Binder oder Drug-Response-Ranking."),
    ("3", "Domaendaten mappen", "Eigene Molekuele, Proteinziele, Zelllinien und Messwerte in MAMMAL-Promptsyntax ueberfuehren."),
    ("4", "Validieren", "Holdout-Sets, negative Kontrollen, Unsicherheit und prospektive Nasslabor-Tests einplanen."),
    ("5", "Entscheiden", "Nur robuste Signale in Priorisierungslisten, Experimentdesigns oder Repurposing-Hypothesen uebernehmen."),
]
for i, (num, head, body) in enumerate(steps):
    y = 1.45 + i * 0.95
    pill(slide, 0.9, y, 0.55, 0.45, num, TEAL if i < 3 else AMBER)
    text_box(slide, 1.65, y + 0.01, 2.0, 0.3, head, 14, True, NAVY)
    text_box(slide, 3.55, y + 0.0, 8.5, 0.4, body, 12.2, False, GRAY)


# 10
slide = prs.slides.add_slide(blank)
add_bg(slide, LIGHT)
title(slide, "Kurzfazit", "Was bleibt fuer Forschung und Produktentscheidungen haengen?")
footer(slide, 10)
text_box(
    slide,
    0.9,
    1.65,
    11.4,
    1.15,
    "MAMMAL ist vor allem ein Integrations- und Priorisierungsmodell: Es verbindet mehrere biomedizinische Modalitaeten, schlaegt spezialisierte Baselines in vielen Benchmarks und liefert verwertbare Signale fuer Folgeexperimente.",
    24,
    True,
    NAVY,
)
card(slide, 1.05, 3.35, 3.35, 1.45, "Kernaussage", "Gemeinsame Sequenzmodellierung kann Drug-Discovery-Aufgaben ueber Modalitaeten hinweg vereinheitlichen.", TEAL)
card(slide, 4.95, 3.35, 3.35, 1.45, "Belege", "9/11 SOTA, AF3-Vergleich 5/7 Targets, wet-lab-bestaetigte Drug-Response-Rangfolge.", GREEN)
card(slide, 8.85, 3.35, 3.35, 1.45, "Entscheidung", "Fuer Exploration und Priorisierung einsetzen; kritische Kandidaten prospektiv validieren.", AMBER)
text_box(slide, 1.1, 5.75, 10.8, 0.45, "Empfohlene erste Anwendung: ein eng definierter Benchmark mit eigenen Daten und klarer experimenteller Nachvalidierung.", 15, False, GRAY, PP_ALIGN.CENTER)


prs.save(OUT)
print(OUT.resolve())
