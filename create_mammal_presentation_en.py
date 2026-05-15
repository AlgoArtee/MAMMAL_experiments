import sys
from pathlib import Path

sys.path.insert(0, str(Path(".deps").resolve()))

from create_mammal_presentation import (  # noqa: E402
    AMBER,
    GREEN,
    GRAY,
    LIGHT,
    MID,
    NAVY,
    OUT as _DE_OUT,
    RED,
    TEAL,
    WHITE,
    add_bg,
    add_table,
    blank,
    bullet_list,
    card,
    footer,
    metric_card,
    pill,
    prs,
    text_box,
    title,
)

from pptx import Presentation  # noqa: E402
from pptx.chart.data import CategoryChartData  # noqa: E402
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Inches  # noqa: E402


OUT = Path("MAMMAL_Key_Findings_Presentation_EN.pptx")


# Reinitialize because importing the German authoring module creates its deck.
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# 1
slide = prs.slides.add_slide(blank)
add_bg(slide, NAVY)
text_box(slide, 0.65, 0.55, 11.7, 0.7, "MAMMAL", 36, True, WHITE)
text_box(
    slide,
    0.68,
    1.34,
    10.6,
    0.7,
    "Molecular Aligned Multi-Modal Architecture and Language for biomedical discovery",
    20,
    False,
    (205, 230, 235) if False else WHITE,
)
slide.shapes[-1].text_frame.paragraphs[0].runs[0].font.color.rgb = __import__("pptx").dml.color.RGBColor(205, 230, 235)
text_box(
    slide,
    0.7,
    2.35,
    8.7,
    1.1,
    "Core message: A multimodal foundation model can use protein, antibody, small-molecule, and gene-expression data in a shared sequence language and reaches new state-of-the-art results on 9 of 11 drug-discovery benchmarks.",
    24,
    True,
    WHITE,
)
metric_card(slide, 0.7, 4.45, "2B", "pretraining samples", TEAL)
metric_card(slide, 3.05, 4.45, "11", "benchmarks", GREEN)
metric_card(slide, 5.4, 4.45, "9/11", "SOTA results", AMBER)
metric_card(slide, 7.75, 4.45, "5/7", "better than AF3", RED)
text_box(slide, 0.7, 6.85, 8.6, 0.22, "Created from s44386-026-00047-4.pdf; original PDF left unchanged.", 8, False, __import__("pptx").dml.color.RGBColor(196, 208, 218))


# 2
slide = prs.slides.add_slide(blank)
add_bg(slide)
title(slide, "Why it matters", "Drug discovery needs models that connect modalities and tasks.")
footer(slide, 2)
card(slide, 0.7, 1.65, 3.75, 1.45, "Starting problem", "Classical pipelines are costly, slow, and face high clinical failure rates; many data types remain modeled in isolation.", RED)
card(slide, 4.8, 1.65, 3.75, 1.45, "Modeling gap", "Molecules, proteins, antibodies, and gene expression are often handled by specialized architectures rather than shared context.", AMBER)
card(slide, 8.9, 1.65, 3.75, 1.45, "Target state", "One model should support prediction, regression, and generation across multiple stages of drug development.", TEAL)
bullet_list(
    slide,
    1.05,
    3.72,
    11.2,
    1.75,
    [
        "The paper positions MAMMAL as an integrative model for target identification, drug response, safety, binding, and antibody design.",
        "The evaluation is deliberately benchmark-driven: existing splits, established metrics, and comparison against specialized SOTA models.",
        "Practical relevance is supported by an experimentally validated drug-ranking result and comparisons with AF3 confidence scores.",
    ],
    15,
)


# 3
slide = prs.slides.add_slide(blank)
add_bg(slide)
title(slide, "Solution: One shared sequence language", "MAMMAL translates heterogeneous biological entities into structured prompts.")
footer(slide, 3)
for x, label, color in [(0.65, "Small molecules\nSMILES", TEAL), (3.25, "Proteins\namino acids", GREEN), (5.85, "Antibodies\nheavy/light chains", AMBER), (8.45, "Gene expression\nranked genes", RED)]:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.65), Inches(2.15), Inches(0.95))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    text_box(slide, x + 0.1, 1.82, 1.95, 0.45, label, 12, True, WHITE, PP_ALIGN.CENTER)
pill(slide, 4.6, 3.05, 3.25, 0.55, "Modular tokenizer + structured prompts", NAVY)
pill(slide, 4.83, 4.05, 2.8, 0.55, "Transformer with encoder-only and encoder-decoder modes", TEAL)
bullet_list(
    slide,
    0.9,
    5.05,
    11.8,
    1.2,
    [
        "Numeric values are not split into text digits; they are projected directly into the embedding space.",
        "The same framework can handle classification, regression, and generation rather than requiring a separate model per task.",
    ],
    14,
)


# 4
slide = prs.slides.add_slide(blank)
add_bg(slide)
title(slide, "Evidence 1: Broad benchmark coverage", "Performance across eleven tasks along the drug-discovery pipeline.")
footer(slide, 4)
rows = [
    ["Benchmark", "Domain", "Metric", "SOTA", "MAMMAL", "Imp."],
    ["Cell type", "GE", "F1", "0.710", "0.763", "+7.5%"],
    ["BBBP", "SM", "AUROC", "0.937", "0.957", "+2.2%"],
    ["ClinTox", "SM", "AUROC", "0.948", "0.986", "+4.0%"],
    ["Cancer Drug 1", "GE+SM", "Pearson", "0.887", "0.917", "+3.4%"],
    ["Cancer Drug 2", "GE+SM", "Pearson", "0.900", "0.931", "+3.4%"],
    ["Ab Infilling", "Protein", "CDRH3-AAR", "0.375", "0.446", "+19.0%"],
    ["PPI ddG", "Protein", "Pearson", "0.663", "0.852", "+28.5%"],
    ["DTI", "Prot.+SM", "NRMSE", "0.942", "0.906", "+3.8%"],
]
add_table(slide, 0.62, 1.45, 7.65, 4.95, rows, [1.65, 0.75, 1.0, 0.75, 0.9, 0.7], 8.8)
chart_data = CategoryChartData()
chart_data.categories = ["Cell", "ClinTox", "Ab Inf.", "PPI", "DTI"]
chart_data.add_series("Relative improvement", [7.5, 4.0, 19.0, 28.5, 3.8])
chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(8.65), Inches(1.7), Inches(3.9), Inches(3.6), chart_data).chart
chart.has_legend = False
chart.value_axis.has_major_gridlines = True
chart.value_axis.tick_labels.font.size = __import__("pptx").util.Pt(8)
chart.category_axis.tick_labels.font.size = __import__("pptx").util.Pt(8)
chart.plots[0].series[0].format.fill.solid()
chart.plots[0].series[0].format.fill.fore_color.rgb = TEAL
text_box(slide, 8.7, 5.58, 3.65, 0.62, "Important: Two tasks are only competitive or narrowly improved; the SOTA claim applies within the selected public benchmarks.", 10.5, False, GRAY)


# 5
slide = prs.slides.add_slide(blank)
add_bg(slide)
title(slide, "Evidence 2: Comparison with AlphaFold 3", "MAMMAL is tested as a sequence model against AF3 confidence scores for binder/non-binder decisions.")
footer(slide, 5)
rows = [
    ["Target", "n", "MAMMAL", "AF3", "Delta"],
    ["HER2 ECD", "60", "0.93", "0.45", "+0.42"],
    ["Albumin", "28", "0.91", "0.59", "+0.32"],
    ["CD206", "37", "1.00", "0.59", "+0.41"],
    ["EGFR", "28", "0.94", "0.49", "+0.45"],
    ["TBG", "32", "0.63", "1.00", "-0.37"],
    ["TNF alpha", "34", "0.86", "0.87", "-0.01"],
    ["VWF", "34", "0.83", "0.32", "+0.51"],
]
add_table(slide, 0.72, 1.45, 5.35, 4.25, rows, [1.3, 0.5, 0.85, 0.65, 0.7], 9.5)
chart_data = CategoryChartData()
chart_data.categories = ["HER2", "Alb.", "CD206", "EGFR", "TBG", "TNF", "VWF"]
chart_data.add_series("MAMMAL", [0.93, 0.91, 1.00, 0.94, 0.63, 0.86, 0.83])
chart_data.add_series("AF3", [0.45, 0.59, 0.59, 0.49, 1.00, 0.87, 0.32])
chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(6.45), Inches(1.48), Inches(6.05), Inches(4.25), chart_data).chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.value_axis.maximum_scale = 1.0
chart.value_axis.tick_labels.font.size = __import__("pptx").util.Pt(8)
chart.category_axis.tick_labels.font.size = __import__("pptx").util.Pt(8)
chart.plots[0].series[0].format.fill.solid()
chart.plots[0].series[0].format.fill.fore_color.rgb = TEAL
chart.plots[0].series[1].format.fill.solid()
chart.plots[0].series[1].format.fill.fore_color.rgb = AMBER
bullet_list(slide, 0.8, 6.05, 11.7, 0.75, ["Interpretation: AF3 provides 3D structural hypotheses; MAMMAL provides discriminative binding probabilities. The comparison is about decision utility, not identical modeling goals."], 12.5)


# 6
slide = prs.slides.add_slide(blank)
add_bg(slide)
title(slide, "Evidence 3: Drug response with wet-lab anchor", "Four drugs absent from training were experimentally tested in cell lines.")
footer(slide, 6)
for i, (drug, note, color) in enumerate([
    ("Carfilzomib", "highest potency", GREEN),
    ("Nintedanib", "rank 2", TEAL),
    ("Infigratinib", "rank 3", AMBER),
    ("Vemurafenib", "lowest potency", RED),
]):
    x = 0.95 + i * 3.0
    arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(1.8), Inches(1.8), Inches(1.05))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    text_box(slide, x - 0.2, 3.05, 2.25, 0.28, drug, 12.5, True, NAVY, PP_ALIGN.CENTER)
    text_box(slide, x - 0.2, 3.38, 2.25, 0.25, note, 9.5, False, GRAY, PP_ALIGN.CENTER)
bullet_list(
    slide,
    1.0,
    4.45,
    11.2,
    1.25,
    [
        "MAMMAL reproduced the experimental potency ranking for the tested cell lines.",
        "Across 805 GDSC cell lines, the relative order was preserved in roughly 90-95% of cases.",
        "Carfilzomib is approved for hematological malignancies; the solid-tumor-cell-line signal supports follow-up, not immediate clinical conclusions.",
    ],
    14,
)


# 7
slide = prs.slides.add_slide(blank)
add_bg(slide)
title(slide, "Decisions implied by the paper", "Which architecture and research decisions carry the results?")
footer(slide, 7)
card(slide, 0.72, 1.55, 3.8, 1.55, "1. Multi-domain over single-domain", "Shared syntax for SMILES, protein sequences, antibody chains, and gene expression.", TEAL)
card(slide, 4.82, 1.55, 3.8, 1.55, "2. Model numbers natively", "Continuous values are projected; this is central for IC50, pKD, ddG, and other scales.", GREEN)
card(slide, 8.92, 1.55, 3.8, 1.55, "3. Two operating modes", "Encoder-only for classification/regression, encoder-decoder for generation and infilling.", AMBER)
card(slide, 0.72, 3.65, 3.8, 1.55, "4. Fair splits for SOTA claims", "Benchmark selection prioritizes public data, defined splits, and established metrics.", NAVY)
card(slide, 4.82, 3.65, 3.8, 1.55, "5. Use AF3 selectively", "AF3 remains valuable for structural hypotheses; MAMMAL is often stronger for binder classification.", RED)
card(slide, 8.92, 3.65, 3.8, 1.55, "6. Release openly", "Code, model weights, and tokenizer are available as a base for reproduction and further studies.", TEAL)


# 8
slide = prs.slides.add_slide(blank)
add_bg(slide)
title(slide, "What not to overinterpret", "The evidence is strong, but not unlimited.")
footer(slide, 8)
bullet_list(
    slide,
    0.9,
    1.55,
    5.75,
    4.5,
    [
        "The SOTA claim depends on the reported public benchmarks and their splits.",
        "AF3 comparisons use confidence scores as proxies; AF3 is not an explicitly trained binder classifier.",
        "Some AF3 tests are small, e.g. HER2 with 60 examples and several targets under 40 examples.",
        "Cell-line and CellTiter-Glo drug-response results are early evidence; clinical transfer remains open.",
    ],
    15,
)
card(slide, 7.2, 1.65, 4.75, 1.45, "Implication", "MAMMAL is a robust research platform for prioritization and hypothesis generation, not a substitute for prospective biological validation.", RED)
card(slide, 7.2, 3.55, 4.75, 1.45, "Governance", "Production use requires dataset-specific validation, uncertainty communication, and explicit task boundaries.", AMBER)


# 9
slide = prs.slides.add_slide(blank)
add_bg(slide)
title(slide, "Next steps", "A pragmatic path from paper to usable decision support.")
footer(slide, 9)
steps = [
    ("1", "Reproduce", "Load code/weights, rebuild target benchmarks locally, and compare metrics with the paper."),
    ("2", "Choose a use case", "Pick one clear workflow: DTI prioritization, antibody binders, or drug-response ranking."),
    ("3", "Map domain data", "Convert internal molecules, protein targets, cell lines, and measurements into MAMMAL prompt syntax."),
    ("4", "Validate", "Plan holdout sets, negative controls, uncertainty checks, and prospective wet-lab tests."),
    ("5", "Decide", "Move only robust signals into prioritization lists, experiment designs, or repurposing hypotheses."),
]
for i, (num, head, body) in enumerate(steps):
    y = 1.45 + i * 0.95
    pill(slide, 0.9, y, 0.55, 0.45, num, TEAL if i < 3 else AMBER)
    text_box(slide, 1.65, y + 0.01, 2.0, 0.3, head, 14, True, NAVY)
    text_box(slide, 3.55, y + 0.0, 8.5, 0.4, body, 12.2, False, GRAY)


# 10
slide = prs.slides.add_slide(blank)
add_bg(slide, LIGHT)
title(slide, "Bottom line", "What matters for research and product decisions?")
footer(slide, 10)
text_box(
    slide,
    0.9,
    1.65,
    11.4,
    1.15,
    "MAMMAL is primarily an integration and prioritization model: it connects multiple biomedical modalities, beats specialized baselines across many benchmarks, and produces actionable signals for follow-up experiments.",
    24,
    True,
    NAVY,
)
card(slide, 1.05, 3.35, 3.35, 1.45, "Core message", "Shared sequence modeling can unify drug-discovery tasks across modalities.", TEAL)
card(slide, 4.95, 3.35, 3.35, 1.45, "Evidence", "9/11 SOTA, AF3 comparison on 5/7 targets, and wet-lab-confirmed drug-response ranking.", GREEN)
card(slide, 8.85, 3.35, 3.35, 1.45, "Decision", "Use it for exploration and prioritization; prospectively validate critical candidates.", AMBER)
text_box(slide, 1.1, 5.75, 10.8, 0.45, "Recommended first application: one narrowly defined benchmark with internal data and explicit experimental follow-up.", 15, False, GRAY, PP_ALIGN.CENTER)


prs.save(OUT)
print(OUT.resolve())
